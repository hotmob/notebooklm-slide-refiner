"""Assemble PPTX deck from images."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from notebooklm_slide_refiner.utils import Resolution

EMU_PER_PIXEL = 9525


def assemble_pptx(image_paths: list[Path], pptx_path: Path, resolution: Resolution) -> None:
    """Create a PPTX with each image as a full-slide background."""

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Emu(resolution.width * EMU_PER_PIXEL)
    presentation.slide_height = Emu(resolution.height * EMU_PER_PIXEL)

    blank_layout = presentation.slide_layouts[6]

    for image_path in image_paths:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            left=Emu(0),
            top=Emu(0),
            width=Emu(resolution.width * EMU_PER_PIXEL),
            height=Emu(resolution.height * EMU_PER_PIXEL),
        )

    presentation.save(pptx_path)
